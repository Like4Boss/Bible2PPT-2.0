import json
import os
import re
from tkinter import Tk, Label, Entry, Button, StringVar, OptionMenu, messagebox
from tkinter.filedialog import askopenfilename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from utils.book_map import book_map, ALL_BOOK_NAMES

# === Load Slide Formatting ===
with open(os.path.join("template", "slide_format.json"), encoding="utf-8") as f:
    slide_style = json.load(f)["default"]

# === Load the Bible JSON ===
with open("bible_data/bible_combined.json", "r", encoding="utf-8") as f:
    bible = json.load(f)

target_ppt_path = None

# === Parse user input ===
def parse_reference(ref_string, bible):
    # Normalize all dashes
    ref_string = re.sub(r"[–—−‑]", "-", ref_string)

    # Pattern to detect: book chapter:verse OR book chapter:verse-verse OR book chapter:verse-chapter:verse
    pattern = r"([^\d\s:]+)\s+(\d+):(\d+)(?:-(\d+)(?::(\d+))?)?"

    matches = re.findall(pattern, ref_string)
    results = []

    for match in matches:
        book, start_ch, start_v, end_ch_or_v, end_v = match
        start_ch, start_v = int(start_ch), int(start_v)

        if not end_ch_or_v:
            # Single verse (e.g., 로마서 1:1)
            results.append((book, str(start_ch), str(start_v)))

        elif end_ch_or_v and not end_v:
            # Same chapter verse range (e.g., 창세기 1:1-3)
            end_v = int(end_ch_or_v)
            for v in range(start_v, end_v + 1):
                results.append((book, str(start_ch), str(v)))

        elif end_ch_or_v and end_v:
            # Cross-chapter range (e.g., 이사야 1:1-2:3)
            end_ch, end_v = int(end_ch_or_v), int(end_v)

            for ch in range(start_ch, end_ch + 1):
                first_verse = start_v if ch == start_ch else 1
                last_verse = end_v if ch == end_ch else len(bible[book][str(ch)])
                for v in range(first_verse, last_verse + 1):
                    results.append((book, str(ch), str(v)))

    return results

def generate_ppt():
    ref = entry.get()
    top_lang = top_language.get()
    bottom_lang = bottom_language.get()

    try:
        verse_refs = parse_reference(ref, bible)
    except Exception as e:
        messagebox.showerror("Error", f"Invalid format. Try: 창 1 or 롬 1-3 or 사 1:3-3:9\n\nDetails: {str(e)}")
        return

    if not verse_refs:
        messagebox.showerror("Error", "No valid verses found.")
        return

    if not target_ppt_path:
        messagebox.showerror("Error", "No PowerPoint file selected.")
        return

    try:
        prs = Presentation(target_ppt_path)
        count = 0

        for book, chapter, verse in verse_refs:
            try:
                top_text = bible[book][chapter][verse][top_lang]
                bottom_text = bible[book][chapter][verse][bottom_lang]
            except KeyError:
                continue  # skip missing verses

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            left = Inches(1.0)
            top_inch = Inches(2.0)
            width = Inches(11.33)
            height = Inches(4.5)

            textbox = slide.shapes.add_textbox(left, top_inch, width, height)
            tf = textbox.text_frame
            tf.clear()

            verse_ref = f"{book} {chapter}:{verse}"
            eng_book_ref = f"{book_map.get(book, book)} {chapter}:{verse}"

            # Top (Korean or primary language)
            top_par = tf.paragraphs[0]
            run1a = top_par.add_run()
            run1a.text = f"{verse_ref} "
            run1a.font.bold = True
            run1a.font.size = Pt(slide_style["font_size"])
            run1a.font.name = slide_style["font_name"]

            run1b = top_par.add_run()
            run1b.text = top_text
            run1b.font.bold = False
            run1b.font.size = Pt(slide_style["font_size"])
            run1b.font.name = slide_style["font_name"]

            # Bottom (English or secondary language)
            bot_par = tf.add_paragraph()
            run2a = bot_par.add_run()
            run2a.text = f"{eng_book_ref} "
            run2a.font.bold = True
            run2a.font.size = Pt(slide_style["font_size"] - 2)
            run2a.font.name = slide_style["font_name"]

            run2b = bot_par.add_run()
            run2b.text = bottom_text
            run2b.font.bold = False
            run2b.font.size = Pt(slide_style["font_size"] - 2)
            run2b.font.name = slide_style["font_name"]

            bot_par.space_before = Pt(12)

            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            count += 1

        prs.save(target_ppt_path)
        messagebox.showinfo("Success", f"{count} slide(s) saved to: {target_ppt_path}")

    except KeyError as e:
        messagebox.showerror("Error", f"Verse not found: {e}")
        return

def choose_ppt_file():
    global target_ppt_path
    path = askopenfilename(
        title="Select PowerPoint File to Insert Slides",
        filetypes=[("PowerPoint Files", "*.pptx")]
    )
    if path:
        target_ppt_path = path
        messagebox.showinfo("Selected", f"Target file set to:\n{path}")

def launch_gui():
    global entry, top_language, bottom_language

    root = Tk()
    root.title("Bible2PPT (KKJV/NIV/NGAYOK)")
    root.geometry("700x400")

    Label(root, text="Enter Verse (e.g. 창세기 1:1 or 요한복음 3:16-18):").pack(pady=5)
    entry = Entry(root, width=40)
    entry.pack()

    Label(root, text="Top Language:").pack()
    top_language = StringVar(root)
    top_language.set("ngayok")
    OptionMenu(root, top_language, "kkjv", "niv", "ngayok").pack()

    Label(root, text="Bottom Language:").pack()
    bottom_language = StringVar(root)
    bottom_language.set("niv")
    OptionMenu(root, bottom_language, "kkjv", "niv", "ngayok").pack()

    Button(root, text="Select PowerPoint File", command=choose_ppt_file).pack(pady=5)
    Button(root, text="Generate PowerPoint", command=generate_ppt).pack(pady=20)

    root.mainloop()

if __name__ == "__main__":
    print("🚀 GUI is launching...")
    launch_gui()
